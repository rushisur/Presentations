from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os
import json
from groq import Groq

def check_api_key(api_key=None):
    """Check if GROQ_API_KEY is set"""
    if not api_key:
        api_key = os.environ.get("GROQ_API_KEY")
    if not api_key:
        raise ValueError(
            "Groq API Key is not set! Please provide your API key."
        )
    return api_key

def generate_slide_content(client, topic, num_slides=4, points_per_slide=3):
    """Generate slide content using Groq"""
    try:
        prompt = f"""Create a professional presentation outline about: {topic}
        Follow these requirements:
        - Total slides: {num_slides}
        - Points per slide: {points_per_slide}
        - Structure in this JSON format:
        {{
            "title": "Presentation Title",
            "slides": [
                {{
                    "title": "Slide Title",
                    "content": ["Point 1", "Point 2", "Point 3"]
                }}
            ]
        }}
        Include these elements:
        1. Introduction slide
        2. Key concept slides
        3. Real-world examples
        4. Conclusion slide
        Keep content concise and professional."""
        
        completion = client.chat.completions.create(
            model="mixtral-8x7b-32768",
            messages=[
                {
                    "role": "system",
                    "content": "You are a presentation expert creating structured slide content."
                },
                {
                    "role": "user", 
                    "content": prompt
                }
            ],
            temperature=0.7
        )
        
        response = completion.choices[0].message.content
        return json.loads(response)
            
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON response: {response}") from e
    except Exception as e:
        raise RuntimeError(f"API Error: {str(e)}") from e

def create_title_slide(prs, title):
    """Create a title slide"""
    try:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title
        return slide
    except Exception as e:
        raise RuntimeError(f"Slide creation error: {str(e)}") from e

def create_content_slide(prs, title, content):
    """Create a content slide with bullet points"""
    try:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in content:
            p = tf.add_paragraph()
            p.text = str(item)
            p.level = 0
            
        return slide
    except Exception as e:
        raise RuntimeError(f"Content slide error: {str(e)}") from e

def generate_presentation(client, topic, num_slides=4, points_per_slide=3):
    """
    Generate complete presentation
    
    Args:
        client: Authenticated Groq client
        topic (str): Presentation topic
        num_slides (int): Number of slides
        points_per_slide (int): Points per slide
    """
    try:
        # Validate inputs
        if num_slides < 1 or points_per_slide < 1:
            raise ValueError("Invalid slide configuration")
            
        # Get AI-generated content
        content = generate_slide_content(
            client=client,
            topic=topic,
            num_slides=num_slides,
            points_per_slide=points_per_slide
        )
        
        # Create presentation
        prs = Presentation()
        create_title_slide(prs, content["title"])
        
        # Validate slide structure
        if len(content["slides"]) != num_slides:
            raise ValueError("Generated slides count mismatch")
            
        for slide in content["slides"]:
            create_content_slide(
                prs,
                slide["title"],
                slide["content"][:points_per_slide]  # Enforce points limit
            )
            
        return prs
        
    except KeyError as e:
        raise ValueError(f"Invalid content format: {str(e)}") from e
    except Exception as e:
        raise RuntimeError(f"Generation failed: {str(e)}") from e